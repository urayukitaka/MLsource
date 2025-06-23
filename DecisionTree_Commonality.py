import os
import sys
sys.path.append(__file__)
from typing import Union
import pandas as pd
import numpy as np
import re
from sklearn.model_selection import train_test_split
from sklearn.impute import SimpleImputer
from sklearn.preprocessing import OneHotEncoder
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline
from sklearn.tree import DecisionTreeClassifier
from sklearn.metrics import accuracy_score, precision_score, recall_score, f1_score
import seaborn as sns
from sklearn.metrics import confusion_matrix
from sklearn.tree import plot_tree, _tree
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
import matplotlib.cm as cm
import matplotlib.colors as mcolors

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

import utils as u

class Commonality:

    def __init__(self,
                 savedir:str):
        print("")
        print("#"*50)
        print("# Commonality class, by Learning of shallow decision trees")
        print("#"*50)
        print("")
        print(" Preparing savedirs. model, report, graph dirs")
        self.savedir = savedir
        os.makedirs(self.savedir, exist_ok=True)
        self.modeldir = os.path.join(self.savedir, "model")
        os.makedirs(self.modeldir, exist_ok=True)
        self.reportdir = os.path.join(self.savedir, "report")
        os.makedirs(self.reportdir, exist_ok=True)
        self.graphdir = os.path.join(self.reportdir, "graphs")
        os.makedirs(self.graphdir, exist_ok=True)

    def pre_analyze_dataset(self,
                            target_col:str,
                            y_data:Union[list, np.array]=[],
                            y_data_numeric:Union[list, np.array]=[],
                            numerical_parametere:list=[],
                            categorical_parameter:list=[]
                            ) -> None:
        print(" Analyze dataset.")
        # dtype
        y_data = np.array(y_data)
        y_data_numeric = np.array(y_data_numeric)

        # check
        self.dataset_info = {}
        self.dataset_info["Data_count"] = len(y_data)
        self.dataset_info["Target_value"] = {
            "Target_name":target_col,
            "Good":{"Count":len(y_data[y_data==0]),
                    "Value_average":np.nan if len(y_data_numeric)==0 else np.mean(y_data_numeric[y_data==0])},
            "Bad":{"Count":len(y_data[y_data==1]),
                   "Value_average":np.nan if len(y_data_numeric)==0 else np.mean(y_data_numeric[y_data==1])}
        }
        self.dataset_info["Explanatory_value"] = {
            "Numeric_params_count":len(numerical_parametere),
            "Categorical_params_count":len(categorical_parameter),
        }

        # make value hist
        if len(y_data_numeric)!=0:
            self.target_val_dist = os.path.join(self.graphdir, "target_val_dist.png")
            # make graph
            boxdf = pd.DataFrame({
                "Good_Bad":["Good" if v == 0 else "Bad" for v in y_data],
                target_col:y_data_numeric
            })
            box_hist_with_stats_colored_auto(df=boxdf, category_col="Good_Bad", value_col=target_col,
                                             figsize=(10,6), savepath=self.target_val_dist)

        else:
            self.target_val_dist = None
        self.dataset_info["Target_value_graph"] = self.target_val_dist

        # save
        self.dataset_info_file = os.path.join(self.modeldir, "dataset_info.json")
        u.save_jsonf(dict_data=self.dataset_info,
                     savefilename=self.dataset_info_file,
                     encoding="CP932")

    def learning(self,
                 dataset:pd.DataFrame,
                 target_col:str,
                 threshold:float,
                 bad_direction:str,
                 numeric:bool,
                 goodlbl:str=None,
                 badlbl:str=None,
                 val_size:float=0.2) -> None:
        # init
        self.dataset = dataset
        self.target_col = target_col
        self.threshold = threshold
        self.numeric = numeric
        self.goodlbl, self.badlbl = goodlbl, badlbl
        self.val_size = val_size

        # initilize learning result
        self.learning_result = {}
        # ------------------------------------------------------
        # Dataset
        # ------------------------------------------------------
        print("")
        print(" Preparing Dataset.")
        # X, y
        X = dataset.drop(columns=target_col)
        if numeric:
            if bad_direction=="under":
                y = dataset[target_col].apply(lambda x:1 if x < threshold else 0)
            elif bad_direction=="over":
                y = dataset[target_col].apply(lambda x:1 if x > threshold else 0)
        else:
            if badlbl!=None:
                y = dataset[target_col].apply(lambda x:1 if x == badlbl else 0)
            else:
                y = dataset[target_col]
        # label y
        self.label_y = y.copy()

        # analyze features
        self.categorical_cols = X.select_dtypes(include=["category", "object"]).columns.tolist()
        self.numeric_cols = X.select_dtypes(include=["number"]).columns.tolist()

        # analyze dataset
        if numeric:
            self.pre_analyze_dataset(
                target_col=target_col,
                y_data=list(y),
                y_data_numeric=dataset[target_col].to_list(),
                numerical_parametere=self.numeric_cols,
                categorical_parameter=self.categorical_cols
            )
        else:
            self.pre_analyze_dataset(
                target_col=target_col,
                y_data=y,
                numerical_parametere=self.numeric_cols,
                categorical_parameter=self.categorical_cols
            )
        print(" -> Done")

        # ------------------------------------------------------
        # Modeling
        # ------------------------------------------------------
        print("")
        print(" Difine model and preprocessing.")
        # make pipeline, preprocessing and model
        num_transformer = Pipeline(
            steps = [
                ("imputer", SimpleImputer(strategy="mean"))
            ]
        )
        cat_transformer = Pipeline(
            steps = [
                ("imputer", SimpleImputer(strategy="constant", fill_value="Missing")),
                ("encoder", OneHotEncoder(handle_unknown="ignore"))
            ]
        )
        preprocessor = ColumnTransformer(
            transformers=[
                ("num", num_transformer, self.numeric_cols),
                ("cat", cat_transformer, self.categorical_cols)
            ]
        )
        self.model = Pipeline(steps=[
            ("preprocess", preprocessor),
            ("clf", DecisionTreeClassifier(max_depth=4, random_state=42))
        ])
        # train test split
        X_train, X_test, y_train, y_test = train_test_split(
            X, y, test_size=val_size, random_state=42, stratify=y
        )

        print(" -> Done")

        # ------------------------------------------------------
        # Learning
        # ------------------------------------------------------
        print("")
        print(" Start learning.")
        # fit
        self.model.fit(X_train, y_train)

        # savemodel
        self.preprocess_file = os.path.join(self.modeldir, "preprocess.pkl")
        u.save_pkl(self.model.named_steps["preprocess"],
                   savefilename=self.preprocess_file)
        self.learning_result["preprocessing_file"] = self.preprocess_file

        self.model_file = os.path.join(self.modeldir, "model.pkl")
        u.save_pkl(self.model.named_steps["clf"],
                   savefilename=self.model_file)
        self.learning_result["model_file"] = self.model_file

        # pipeline
        self.pipeline_file = os.path.join(self.modeldir, "pipeline.pkl")
        u.save_pkl(self.model,
                   savefilename=self.pipeline_file)
        self.learning_result["pipeline_file"] = self.pipeline_file

        self.learnig_result_file = os.path.join(self.modeldir, "learning_result.json")
        u.save_jsonf(self.learning_result,
                     self.learnig_result_file,
                     encoding="CP932"
                     )
        print(" -> Done")

        # ------------------------------------------------------
        # Evaluation
        # ------------------------------------------------------
        print("")
        print(" Evaluate model.")
        self.evals = {}
        y_pred_test = self.model.predict(X_test)
        y_pred_train = self.model.predict(X_train)
        self.evals["Validation"] = {
            "Data_count":len(y_test),
            "Accuracy":accuracy_score(y_test, y_pred_test),
            "Recall":recall_score(y_test, y_pred_test, zero_division=0),
            "Precision":precision_score(y_test, y_pred_test, zero_division=0),
            "F1_score":f1_score(y_test, y_pred_test, zero_division=0),
        }
        self.evals["Train"] = {
            "Data_count":len(y_train),
            "Accuracy":accuracy_score(y_train, y_pred_train),
            "Recall":recall_score(y_train, y_pred_train, zero_division=0),
            "Precision":precision_score(y_train, y_pred_train, zero_division=0),
            "F1_score":f1_score(y_train, y_pred_train, zero_division=0),
        }
        # make confmat
        confmat_file = os.path.join(self.modeldir, "model_confmat.png")
        saveconfmat(y_true=y_test,
                    y_pred=y_pred_test,
                    label_0="<={}".format(threshold),
                    label_1=">{}".format(threshold),
                    savefilename=confmat_file)
        self.evals["ConfusionMatrix"] = confmat_file

        self.eval_file = os.path.join(self.modeldir, "evals.json")
        u.save_jsonf(self.evals,
                     self.eval_file,
                     encoding="CP932"
                     )
        print(" -> Done")

    def analyze(self) -> None:

        # initilize analyze result
        self.analyze_result = {}

        print("")
        print(" Analyze learned model.")

        # *******************
        # (1) tree image
        tree_model = self.model.named_steps["clf"]
        cat_pipeline = self.model.named_steps["preprocess"].named_transformers_["cat"]
        encoder = cat_pipeline.named_steps["encoder"]

        # ③ 特徴量名の復元
        encoded_cat_features = encoder.get_feature_names_out(self.categorical_cols)
        feature_names = self.numeric_cols + list(encoded_cat_features)

        # プロット
        self.savetree = os.path.join(self.graphdir, "tree_structure.png")

        plt.figure(figsize=(14, 8))
        if self.numeric==True:
            plot_tree(
                tree_model,
                feature_names=feature_names,
                class_names=["<={}".format(self.threshold), ">{}".format(self.threshold)],
                filled=True,
                rounded=True,
                fontsize=8
            )
        else:
            if self.badlbl!=None:
                plot_tree(
                    tree_model,
                    feature_names=feature_names,
                    class_names=[self.goodlbl, self.badlbl],
                    filled=True,
                    rounded=True,
                    fontsize=8
                )
            else:
                plot_tree(
                    tree_model,
                    feature_names=feature_names,
                    class_names=["Good", "Bad"],
                    filled=True,
                    rounded=True,
                    fontsize=8
                )
        plt.savefig(self.savetree, bbox_inches='tight')
        plt.close()

        self.analyze_result["tree_graph"] = self.savetree

        # *******************
        # importance
        # 重要度取得（ジニ不純度の減少量に基づく）
        importances = tree_model.feature_importances_

        # データフレームでソート表示
        importance_df = pd.DataFrame({
            "feature": feature_names,
            "importance": importances
        }).sort_values("importance", ascending=False)
        # save
        self.importance_df_path = os.path.join(
            self.reportdir, "importance_result.csv"
        )
        importance_df.to_csv(self.importance_df_path, index=False, encoding="CP932")
        self.analyze_result["importance_df"] = self.importance_df_path

        # graph
        self.importance_graph = os.path.join(
            self.graphdir, "importance_result.png"
        )
        topn = 10
        plt.figure(figsize=(12, 6))
        plt.barh(importance_df["feature"].head(topn)[::-1], importance_df["importance"].head(topn)[::-1])
        plt.xlabel("Gini Importance")
        plt.title(f"Top {topn} Feature Importances (Gini)")
        plt.grid()
        plt.tight_layout()
        plt.savefig(self.importance_graph)
        plt.close()
        self.analyze_result["importance_graph"] = self.importance_graph

        # top3 histgrams
        features = importance_df["feature"].values
        for i in range(6):
            ranking = i + 1
            param = features[i]
            try:
                #if len(param.split("_"))>1:
                if re.match(r"^([^\_]+)", param).group(1):
                    param_col = param.split("_")[0]
            except:
                param_col = param
            # box hist
            histdf = pd.DataFrame({
                param_col:self.dataset[param_col],
                self.target_col:self.dataset[self.target_col]
            })
            savefilename = os.path.join(self.graphdir, "boxplot_importance_rank{}.png".format(ranking))
            box_hist_with_stats_colored_auto(
                df = histdf, category_col = param_col, value_col = self.target_col,
                figsize=(16,6), savepath = savefilename
            )
            self.analyze_result["importance_rank{}_graph".format(ranking)] = savefilename

        # *******************
        # node graphs, by 3 layer
        tree_struct = get_tree_structure_with_paths(tree_model, feature_names, max_depth=4)
        tree_converted = u.convert_numpy(tree_struct)
        self.tree_structure_file = os.path.join(self.reportdir, "tree_structure.json")
        u.save_jsonf(tree_converted,
                     self.tree_structure_file, encoding="CP932")
        self.branch = {}
        # 1st branch
        branch1_filename = os.path.join(self.graphdir, "branch1.png")
        param, condition1, condition1_else, bratio_str, bratio_l_str, bratio_r_str, data_l, data_r = analyze_node(tree_structure=tree_struct,
                                                                                     dataset = self.dataset,
                                                                                     savefilename=branch1_filename)
        self.branch["top"] = {
            "param":param,
            "condition_left":condition1,
            "condition_right":condition1_else,
            "base_bad_ratio":bratio_str,
            "bad_ratio_left":bratio_l_str,
            "bad_ratio_right":bratio_r_str,
            "branch_graph":branch1_filename
        }
        # 2nd left
        branch1_left_filename = os.path.join(self.graphdir, "branch1_left.png")
        param, condition2, condition2_else, _, bratio_l_l_str, bratio_l_r_str, _, _ = analyze_node(tree_structure=tree_struct["left"],
                                                                                dataset = data_l,
                                                                                savefilename=branch1_left_filename)
        self.branch["left"] = {
            "param":param,
            "condition_left":condition2,
            "condition_right":condition2_else,
            "bad_ratio_left":bratio_l_l_str,
            "bad_ratio_right":bratio_l_r_str,
            "branch_graph":branch1_left_filename
        }
        # 3rd right
        branch1_right_filename = os.path.join(self.graphdir, "branch1_right.png")
        param, condition3, condition3_else, _, bratio_r_l_str, bratio_r_r_str, _, _ = analyze_node(tree_structure=tree_struct["right"],
                                                                                dataset = data_r,
                                                                                savefilename=branch1_right_filename)
        self.branch["right"] = {
            "param":param,
            "condition_left":condition3,
            "condition_right":condition3_else,
            "bad_ratio_left":bratio_r_l_str,
            "bad_ratio_right":bratio_r_r_str,
            "branch_graph":branch1_right_filename
        }
        self.branch_file = os.path.join(self.reportdir, "branch.json")
        u.save_jsonf(self.branch,
                     self.branch_file, encoding="CP932")

        # *******************
        # save analysis result
        self.analyze_result["tree_structure"] = self.tree_structure_file
        self.analyze_result_file = os.path.join(self.reportdir, "analyze_result.json")
        u.save_jsonf(self.analyze_result,
                     self.analyze_result_file, encoding="CP932")

def saveconfmat(y_true:np.array,
                y_pred:np.array,
                label_0:str,
                label_1:str,
                savefilename:str) -> None:
    cm = confusion_matrix(y_true, y_pred)
    labels = [label_0, label_1]

    plt.figure(figsize=(6, 5))
    sns.heatmap(cm, annot=True, fmt="d", cmap="Blues",
                xticklabels=labels, yticklabels=labels)
    plt.xlabel("Predicted")
    plt.ylabel("Actual")
    plt.title("Confusion Matrix")
    plt.tight_layout()
    plt.savefig(savefilename)
    plt.close()

def box_hist_with_stats_colored_auto(df: pd.DataFrame,
                                     category_col: str,
                                     value_col: str,
                                     figsize=(10, 6),
                                     savepath=None) -> None:

    if np.issubdtype(df[category_col].dtype, np.number):
        df = df.copy()
        df["category"] = pd.cut(df[category_col], bins=8)
        working_col = "category"
    else:
        working_col = category_col

    unique_vals = df[working_col].unique()
    if isinstance(unique_vals[0], pd.Interval):
        categories = sorted(unique_vals, key=lambda x: x.left)
    else:
        categories = sorted(unique_vals, key=lambda x: str(x))

    x_pos = np.arange(len(categories))
    stats = []

    # カラーマップで色自動割り当て
    cmap = cm.get_cmap('tab10', len(categories))
    colors = [mcolors.to_hex(cmap(i)) for i in range(len(categories))]

    fig, ax = plt.subplots(figsize=figsize)

    # ---- 箱ひげ図（背景にするためzorderを低く） ----
    for i, cat in enumerate(categories):
        subset = df[df[working_col] == cat][value_col].dropna()

        ax.boxplot(
            subset,
            positions=[x_pos[i]],
            widths=0.3,
            patch_artist=True,
            boxprops=dict(facecolor=colors[i], color=colors[i], zorder=1),
            medianprops=dict(color="white", zorder=1),
            whiskerprops=dict(color=colors[i], zorder=1),
            capprops=dict(color=colors[i], zorder=1),
            flierprops=dict(marker='o', color=colors[i], alpha=0.5, zorder=1)
        )

        # 統計値収集
        stats.append([
            np.round(subset.mean(), 2),
            np.round(subset.median(), 2),
            np.round(subset.max(), 2),
            np.round(subset.min(), 2),
            len(subset)
        ])

    # ---- ヒストグラム（前面にするためzorderを高く） ----
    for i, cat in enumerate(categories):
        subset = df[df[working_col] == cat][value_col].dropna()

        hist_vals, bins = np.histogram(subset, bins=20, density=True)
        bin_centers = 0.5 * (bins[1:] + bins[:-1])
        scale = 0.4 / hist_vals.max() if hist_vals.max() > 0 else 1

        for h, y in zip(hist_vals, bin_centers):
            rect = Rectangle(
                (x_pos[i] - h * scale / 2, y),
                h * scale,
                bins[1] - bins[0],
                alpha=0.4,
                color=colors[i],
                zorder=2  # 前面
            )
            ax.add_patch(rect)

    # 軸設定
    ax.set_xticks(x_pos)
    ax.set_xticklabels([str(c) for c in categories], rotation=0, ha='right')
    ax.set_ylabel(value_col)
    ax.set_title(f"{category_col} ごとの分布", fontsize=14)
    ax.grid()

    # ---- 統計テーブル ----
    stat_labels = ["平均", "中央値", "最大", "最小", "件数"]
    table_data = list(zip(*stats))
    cell_text = [[str(val) for val in row] for row in table_data]

    table = plt.table(cellText=cell_text,
                      rowLabels=stat_labels,
                      colLabels=categories,
                      cellLoc='center',
                      rowLoc='center',
                      loc='bottom',
                      bbox=[0.0, -0.45, 1, 0.3])
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    for key, cell in table.get_celld().items():
        cell.set_linewidth(0.5)
        cell.set_edgecolor("gray")

    plt.subplots_adjust(bottom=0.35)

    if savepath:
        plt.savefig(savepath, bbox_inches='tight', dpi=300)
    plt.close()

def get_tree_structure_with_paths(tree_model, feature_names, max_depth=5):
    """
    決定木のノードごとに階層・経路情報を含めた構造辞書を生成
    """
    tree_ = tree_model.tree_

    def recurse(node_id, depth, path_list, path_id):
        is_leaf = tree_.feature[node_id] == _tree.TREE_UNDEFINED
        value = tree_.value[node_id][0]
        node_info = {
            "node_id": node_id,
            "depth": depth,
            "path": path_list,
            "path_id": path_id,
            "samples": int(np.sum(value)),
            "class_distribution": value.tolist(),
            "predicted_class": int(np.argmax(value)),
            "leaf": is_leaf,
        }

        if not is_leaf and depth < max_depth:
            feature = feature_names[tree_.feature[node_id]]
            threshold = tree_.threshold[node_id]
            node_info.update({
                "condition": f"{feature} <= {threshold:.3f}",
                "feature": feature,
                "threshold": threshold,
                "left": recurse(
                    tree_.children_left[node_id],
                    depth + 1,
                    path_list + ["left"],
                    f"{path_id}-0"
                ),
                "right": recurse(
                    tree_.children_right[node_id],
                    depth + 1,
                    path_list + ["right"],
                    f"{path_id}-1"
                ),
            })
        return node_info

    return recurse(0, 0, ["root"], "0")

def out_feat_val(val:str):
    # split
    val_split = val.split("_")
    size = len(val_split)
    # feat
    feat = ""
    for i in range(len(val_split)):
        if i==size-1:
            pass
        elif i!=len(val_split)-2:
            feat += val_split[i] + "_"
        else:
            feat += val_split[i]
    # val
    val = val_split[-1]
    return feat, val

def analyze_node(tree_structure:dict,
                 dataset:pd.DataFrame,
                 savefilename:str):

    # 1st parameter
    param1 = tree_structure["feature"]
    if tree_structure["threshold"]!=0.5:
        param = param1
        # numeric value
        thre = tree_structure["threshold"]
        # left side
        condition1 = tree_structure["condition"]
        condition2 = "Else"
        histdf1 = pd.DataFrame({
            param1:["0:"+condition1 if v<thre else "1:Else" for v in dataset[param]],
            target_col:dataset[target_col]
        })
        # data left and bad ratio left
        data_l = dataset[dataset[param]<thre]
        bratio_l = tree_structure["left"]["class_distribution"][1] / tree_structure["left"]["samples"]
        bratio_l_str = "{}%".format(str(round(bratio_l*100, 1)))
        # data right and bad ratio right
        data_r = dataset[dataset[param]>=thre]
        bratio_r = tree_structure["right"]["class_distribution"][1] / tree_structure["right"]["samples"]
        bratio_r_str = "{}%".format(str(round(bratio_r*100, 1)))
        # base bratio
        bratio_base = (tree_structure["left"]["class_distribution"][1]+tree_structure["right"]["class_distribution"][1]) / (tree_structure["left"]["samples"]+tree_structure["right"]["samples"])
        bratio_base_str = "{}%".format(str(round(bratio_base*100, 1)))

    else:
        # categorical value
        param, val = out_feat_val(tree_structure["feature"])
        # left side
        condition1 = "Else " + tree_structure["feature"]
        condition2 = tree_structure["feature"]
        histdf1 = pd.DataFrame({
            param1:["0:Else {}".format(val)+condition1 if v!=val else f"1:{val}" for v in dataset[param]],
            target_col:dataset[target_col]
        })
        # data right and bad ratio 2
        data_l = dataset[dataset[param]!=val]
        bratio_l = tree_structure["left"]["class_distribution"][1] / tree_structure["left"]["samples"]
        bratio_l_str = "{}%".format(str(round(bratio_l*100, 1)))
        # data right and bad ratio right
        data_r = dataset[dataset[param]==val]
        bratio_r = tree_structure["right"]["class_distribution"][1] / tree_structure["right"]["samples"]
        bratio_r_str = "{}%".format(str(round(bratio_r*100, 1)))
        # base bratio
        bratio_base = (tree_structure["left"]["class_distribution"][1]+tree_structure["right"]["class_distribution"][1]) / (tree_structure["left"]["samples"]+tree_structure["right"]["samples"])
        bratio_base_str = "{}%".format(str(round(bratio_base*100, 1)))

    # make graph
    box_hist_with_stats_colored_auto(
            df = histdf1, category_col = param1, value_col = target_col,
            figsize=(16,6), savepath = savefilename
        )

    # return
    return param, condition1, condition2, bratio_base_str, bratio_l_str, bratio_r_str, data_l, data_r

class Report:

    def __init__(self,
                 savefilename:str,
                 dataset_info:dict,
                 learning_result:dict,
                 evalation_result:dict,
                 analyze_result:dict,
                 branch_result:dict
                 ):
        print("")
        print("#"*50)
        print("# Report class, from learning and analyze result")
        print("#"*50)
        print("")
        # init args
        self.savefilename = savefilename
        self.dataset_info = dataset_info
        self.learning_result = learning_result
        self.evalation_result = evalation_result
        self.analyze_result = analyze_result
        self.branch = branch_result
        # init slide
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)  # 16:9
        self.prs.slide_height = Inches(7.5)
        self.blank_slide_layout = self.prs.slide_layouts[6]
    # --------------------------------
    def make_report(self) -> None:
        # 1st, slide title

        # 2nd, dataset difinition
        self.report_data_difinition()

        # 3rd, learning result
        self.report_learning_result()

        # 4th, analyze result
        self.report_analyze_result()

        # 5th, branch
        self.report_with_diagram()

        # 6th, tree structure
        self.report_decision_tree_structure()

    def report_data_difinition(self) -> None:
        print("")
        print(" Data difinitin slide.")
        # color setting
        navy = RGBColor(0, 32, 96)  # 濃い紺色
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        # --------------------------------
        # slide and title
        # --------------------------------
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        topcomment = "データ定義"
        subcomment = "解析に使用したデータセットの定義は以下"
        self._add_title(slide, topcomment, subcomment)
        # --------------------------------
        # left side, matrix of data difinition
        # --------------------------------
        top_margin = Inches(1.5)  # タイトル下
        # Title
        titlebox = slide.shapes.add_textbox(Inches(0.7), top_margin, Inches(3), Inches(0.3))
        titlebox.text_frame.text = "対象のデータ"
        titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
        titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        # Matrix
        # make left matrix
        left_matrix = self.dataset_info_to_matrix()
        table_left = Inches(1.0)
        table_top = Inches(2.)
        table_width_in = 5.5
        table_width = Inches(table_width_in)
        rows, cols = left_matrix.shape
        row_height_in = 0.3
        table_height_in = Inches((rows + 1) * row_height_in)
        table = slide.shapes.add_table(
            rows + 1,
            cols,
            table_left,
            table_top,
            table_width,
            table_height_in,
        ).table
        width_ratio = [0.65, 0.35]
        total_ratio = sum(width_ratio)
        for idx, ratio in enumerate(width_ratio):
            table.columns[idx].width = Inches((ratio / total_ratio) * table_width_in)
        # header row
        for c in range(cols):
            cell = table.cell(0, c)
            cell.text = str(left_matrix.columns[c])
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "メイリオ"
            run.font.bold = True
            run.font.size = Pt(9)
            run.font.color.rgb = white
            cell.fill.solid()
            cell.fill.fore_color.rgb = navy

        # data rows
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r + 1, c)
                cell.text = str(left_matrix.iat[r, c])
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.name = "メイリオ"
                run.font.size = Pt(9)
                run.font.color.rgb = black
                cell.fill.solid()
                cell.fill.fore_color.rgb = white

        # --------------------------------
        # right side, boxplot graph
        # --------------------------------
        left_margin = Inches(7.)
        graph_w = Inches(5.5)
        graph_h = Inches(3.5)
        # Title
        titlebox = slide.shapes.add_textbox(left_margin, top_margin, Inches(3), Inches(0.3))
        titlebox.text_frame.text = "目的変数の分布, 連続変数の指定あれば"
        titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
        titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        # add graph
        self._add_picture_if_exists(slide,
                                    path=self.dataset_info["Target_value_graph"],
                                    left = left_margin,
                                    top = table_top,
                                    width=graph_w,
                                    height=graph_h)
        self.save_pptx()
        print(" -> Done")

    def report_learning_result(self) -> None:
        print("")
        print(" Learning result slide.")
        # color setting
        darkgreen = RGBColor(0,100,0) # 濃い緑
        darkred = RGBColor(139, 0, 0)  # 濃い赤色
        lightpink = RGBColor(255, 182, 193)  # 薄いピンク
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        # --------------------------------
        # slide and title
        # --------------------------------
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        topcomment = "学習結果"
        subcomment = "決定木モデル学習の結果は以下"
        self._add_title(slide, topcomment, subcomment)
        # --------------------------------
        # left side, matrix of evaluation
        # --------------------------------
        top_margin = Inches(1.5)  # タイトル下
        # Title
        titlebox = slide.shapes.add_textbox(Inches(0.7), top_margin, Inches(3), Inches(0.3))
        titlebox.text_frame.text = "各種精度値"
        titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
        titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        # Matrix
        # make left matrix
        left_matrix = self.eval_info_to_matrix()
        table_left = Inches(1.0)
        table_top = Inches(2.)
        table_width_in = 5.5
        table_width = Inches(table_width_in)
        rows, cols = left_matrix.shape
        row_height_in = 0.3
        table_height_in = Inches((rows + 1) * row_height_in)
        table = slide.shapes.add_table(
            rows + 1,
            cols,
            table_left,
            table_top,
            table_width,
            table_height_in,
        ).table
        width_ratio = [0.65, 0.35]
        total_ratio = sum(width_ratio)
        for idx, ratio in enumerate(width_ratio):
            table.columns[idx].width = Inches((ratio / total_ratio) * table_width_in)
        # header row
        for c in range(cols):
            cell = table.cell(0, c)
            cell.text = str(left_matrix.columns[c])
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "メイリオ"
            run.font.bold = True
            run.font.size = Pt(9)
            run.font.color.rgb = white
            cell.fill.solid()
            cell.fill.fore_color.rgb = darkgreen

        # data rows
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r + 1, c)
                cell.text = str(left_matrix.iat[r, c])
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.name = "メイリオ"
                run.font.size = Pt(9)
                run.font.color.rgb = black
                cell.fill.solid()
                cell.fill.fore_color.rgb = white

        # --------------------------------
        # right side, confmat graph
        # --------------------------------
        left_margin = Inches(7.)
        graph_w = Inches(5.5)
        graph_h = Inches(4.5)
        # Title
        titlebox = slide.shapes.add_textbox(left_margin, top_margin, Inches(3), Inches(0.3))
        titlebox.text_frame.text = "評価データに対する混同行列"
        titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
        titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        # add graph
        self._add_picture_if_exists(slide,
                                    path=self.evalation_result["ConfusionMatrix"],
                                    left = left_margin,
                                    top = table_top,
                                    width=graph_w,
                                    height=graph_h)
        self.save_pptx()
        print(" -> Done")

    def report_analyze_result(self)  -> None:
        print("")
        print(" Importance ranking result slide.")
        # color setting
        darkgreen = RGBColor(0,100,0) # 濃い緑
        darkred = RGBColor(139, 0, 0)  # 濃い赤色
        lightpink = RGBColor(255, 182, 193)  # 薄いピンク
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        # 2 slide, ranking 1~3 and 4~6
        # make left matrix
        left_matrix = self.importance_matrix()
        for ranks in ((1,3),(4,6)):
            # --------------------------------
            # slide and title
            # --------------------------------
            slide = self.prs.slides.add_slide(self.blank_slide_layout)
            topcomment = "コモナリティ解析結果, 連続変数の指定あり"
            subcomment = "重要度(ジニ不純度の減少の寄与度の総和)ランキング, {} ~ {}".format(ranks[0], ranks[1])
            self._add_title(slide, topcomment, subcomment)
            # --------------------------------
            # left upper side, matrix of importance
            # --------------------------------
            top_margin = Inches(0.9)  # タイトル下
            # Title
            titlebox = slide.shapes.add_textbox(Inches(0.7), top_margin, Inches(3), Inches(0.3))
            titlebox.text_frame.text = "重要特徴量の上位"
            titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
            titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
            # Matrix
            table_left = Inches(1.0)
            table_top = Inches(1.2)
            table_width_in = 5.5
            table_width = Inches(table_width_in)
            rows, cols = left_matrix.shape
            row_height_in = 0.25
            table_height_in = Inches((rows + 1) * row_height_in)
            table = slide.shapes.add_table(
                rows + 1,
                cols,
                table_left,
                table_top,
                table_width,
                table_height_in,
            ).table
            width_ratio = [0.2, 0.8]
            total_ratio = sum(width_ratio)
            for idx, ratio in enumerate(width_ratio):
                table.columns[idx].width = Inches((ratio / total_ratio) * table_width_in)
            # header row
            for c in range(cols):
                cell = table.cell(0, c)
                cell.text = str(left_matrix.columns[c])
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.name = "メイリオ"
                run.font.bold = True
                run.font.size = Pt(9)
                run.font.color.rgb = white
                cell.fill.solid()
                cell.fill.fore_color.rgb = darkred

            # data rows
            for r in range(rows):
                for c in range(cols):
                    cell = table.cell(r + 1, c)
                    cell.text = str(left_matrix.iat[r, c])
                    run = cell.text_frame.paragraphs[0].runs[0]
                    run.font.name = "メイリオ"
                    run.font.size = Pt(9)
                    run.font.color.rgb = black
                    cell.fill.solid()
                    if r+1 >= ranks[0] and r+1 <= ranks[1]:
                        cell.fill.fore_color.rgb = lightpink
                    else:
                        cell.fill.fore_color.rgb = white

            # --------------------------------
            # left lower side, importance graph
            # --------------------------------
            top_margin2 = top_margin + Inches(2.5)
            graph_w = Inches(5.5)
            graph_h = Inches(3.)
            # Title
            titlebox = slide.shapes.add_textbox(Inches(0.7), top_margin2, Inches(3), Inches(0.3))
            titlebox.text_frame.text = "Importanceグラフ"
            titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
            titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
            # add graph
            self._add_picture_if_exists(slide,
                                        path=self.analyze_result["importance_graph"],
                                        left = table_left,
                                        top = top_margin2 + Inches(0.3),
                                        width=graph_w,
                                        height=graph_h)

            # --------------------------------
            # right side, boxplot graph
            # --------------------------------
            left_margin = Inches(7.)
            graph_w_in = Inches(4.0)
            graph_h = 1.6
            graph_h_in = Inches(graph_h)
            graph_area = graph_h+0.3
            # Title
            titlebox = slide.shapes.add_textbox(left_margin, top_margin, Inches(3), Inches(0.3))
            titlebox.text_frame.text = "ランキング {}～{}, 項目別の分布".format(ranks[0], ranks[1])
            titlebox.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            titlebox.text_frame.paragraphs[0].runs[0].font.bold = True
            titlebox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"

            # add graph
            top_margin2 = top_margin + Inches(0.3)
            for i in range(ranks[0], ranks[1]+1):
                valname = left_matrix["変数"].values[i-1]
                graph_path = self.analyze_result[f"importance_rank{i}_graph"]
                textbox = slide.shapes.add_textbox(left_margin, top_margin2 + Inches( ((i-1)%3) * graph_area), Inches(3), Inches(0.3))
                textbox.text_frame.text = "ランキング {}, {}".format(i, valname)
                textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                textbox.text_frame.paragraphs[0].runs[0].font.bold = True
                textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
                self._add_picture_if_exists(slide,
                                            path=graph_path,
                                            left = left_margin + Inches(1.2),
                                            top = top_margin2 + Inches(0.3) + Inches(((i-1)%3) * graph_area),
                                            width=graph_w_in,
                                            height=graph_h_in)
        self.save_pptx()
        print(" -> Done")

    def report_with_diagram(self)  -> None:
        print("")
        print(" Make diagram slide.")
        # color setting
        darkblue = RGBColor(0,0,200) # 濃い緑
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        # --------------------------------
        # slide and title
        # --------------------------------
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        topcomment = "学習結果"
        subcomment = "3層目までのブランチと各カテゴリの分布"
        self._add_title(slide, topcomment, subcomment)
        # --------------------------------
        # Branch, matrix of evaluation
        # --------------------------------
        c_start_x, c_start_y = Inches(0.7), Inches(1.0)
        xmargin = Inches(0.6)
        text_box_h = Inches(0.3)
        text_box_w_1st = Inches(1.0)
        text_box_w = Inches(2.5)

        # 1st
        root = self.add_box(slide,
                            "ALL",
                            c_start_x,
                            c_start_y,
                            text_box_w_1st,
                            text_box_h,
                            darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x - Inches(0.1),
                                           c_start_y + Inches(0.3), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["top"]["base_bad_ratio"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"

        # 2nd
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + xmargin, c_start_y, Inches(3), Inches(0.3))
        textbox.text_frame.text = "変数 : {}".format(self.branch["top"]["param"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"

        child2 = self.add_box(slide,
                                "{}".format(self.branch["top"]["condition_left"]),
                                c_start_x + text_box_w_1st + xmargin,
                                c_start_y + Inches(0.3),
                                text_box_w,
                                text_box_h,
                                darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + xmargin,
                                           c_start_y + Inches(0.6), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["top"]["bad_ratio_left"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        self.connect(root, child2, slide)

        # 2nd-1st
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                           c_start_y + Inches(0.3), Inches(3), Inches(0.3))
        textbox.text_frame.text = "変数 : {}".format(self.branch["left"]["param"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"

        child21 = self.add_box(slide,
                                "{}".format(self.branch["left"]["condition_left"]),
                                c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                c_start_y + Inches(0.6),
                                text_box_w,
                                text_box_h,
                                darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                           c_start_y + Inches(0.9), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["left"]["bad_ratio_left"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        self.connect(child2, child21, slide)

        # 2nd-2nd
        child22 = self.add_box(slide,
                                "{}".format(self.branch["left"]["condition_right"]),
                                c_start_x + text_box_w_1st + text_box_w + xmargin*2,
                                c_start_y + Inches(2.0),
                                text_box_w,
                                text_box_h,
                                darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                           c_start_y + Inches(2.3), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["left"]["bad_ratio_right"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        self.connect(child2, child22, slide)


        # 3rd
        child3 = self.add_box(slide,
                                "{}".format(self.branch["top"]["condition_right"]),
                                c_start_x + text_box_w_1st + xmargin,
                                c_start_y + Inches(3.2),
                                text_box_w,
                                text_box_h,
                                darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + xmargin,
                                           c_start_y + Inches(3.5), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["top"]["bad_ratio_right"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        self.connect(root, child3, slide)


        # 3rd-1st
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                           c_start_y + Inches(3.2), Inches(3), Inches(0.3))
        textbox.text_frame.text = "変数 : {}".format(self.branch["right"]["param"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"

        child31 = self.add_box(slide,
                                "{}".format(self.branch["right"]["condition_left"]),
                                c_start_x + text_box_w_1st + text_box_w + xmargin*2,
                                c_start_y + Inches(3.5),
                                text_box_w,
                                text_box_h,
                                darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                           c_start_y + Inches(3.8), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["right"]["bad_ratio_left"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        self.connect(child3, child31, slide)

        # 3rd-2nd
        child32 = self.add_box(slide,
                                "{}".format(self.branch["right"]["condition_right"]),
                                c_start_x + text_box_w_1st + text_box_w + xmargin*2,
                                c_start_y + Inches(5.0),
                                text_box_w,
                                text_box_h,
                                darkblue)
        # text box
        textbox = slide.shapes.add_textbox(c_start_x + text_box_w_1st + text_box_w +  xmargin*2,
                                           c_start_y + Inches(5.3), Inches(1.), Inches(0.3))
        textbox.text_frame.text = "Bad比率 : {}".format(self.branch["right"]["bad_ratio_right"])
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
        textbox.text_frame.paragraphs[0].runs[0].font.bold = False
        textbox.text_frame.paragraphs[0].runs[0].font.name = "メイリオ"
        self.connect(child3, child32, slide)

        # add graphs
        # root branch
        graph_w = Inches(4.5)
        graph_h = Inches(2.0)
        self._add_picture_if_exists(slide,
                                    path=self.branch["top"]["branch_graph"],
                                    left=Inches(0.4),
                                    top=Inches(2.0),
                                    width = graph_w,
                                    height = graph_h)
        # left branch
        self._add_picture_if_exists(slide,
                                    path=self.branch["left"]["branch_graph"],
                                    left=c_start_x + text_box_w_1st + text_box_w +  xmargin*2 + text_box_w + Inches(0.1),
                                    top=c_start_y + Inches(0.6),
                                    width = graph_w,
                                    height = graph_h)
        # right branch
        self._add_picture_if_exists(slide,
                                    path=self.branch["right"]["branch_graph"],
                                    left=c_start_x + text_box_w_1st + text_box_w +  xmargin*2 + text_box_w + Inches(0.1),
                                    top=c_start_y + Inches(3.5),
                                    width = graph_w,
                                    height = graph_h)
        self.save_pptx()
        print(" -> Done")

    def report_decision_tree_structure(self)  -> None:
        print("")
        print(" Decision tree slide.")
        # --------------------------------
        # slide and title
        # --------------------------------
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        topcomment = "決定木構造"
        subcomment = "全体構造, 木の最深さ=4"
        self._add_title(slide, topcomment, subcomment)
        # --------------------------------
        # side, decision tree
        # --------------------------------
        top_margin = Inches(1.0)
        left_margin = Inches(0.5)
        graph_w = Inches(12.5)
        graph_h = Inches(6.)
        # add graph
        self._add_picture_if_exists(slide,
                                    path=self.analyze_result["tree_graph"],
                                    left = left_margin,
                                    top = top_margin,
                                    width=graph_w,
                                    height=graph_h)
        self.save_pptx()
        print(" -> Done")

    def _add_title(self, slide, topline: str, subline: str) -> None:
        """Common two‑level title box used by every slide."""
        box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.8))
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True

        run_main = tf.paragraphs[0].add_run()
        run_main.text = topline
        run_main.font.size = Pt(20)
        run_main.font.bold = True
        run_main.font.name = "メイリオ"

        run_sub = tf.add_paragraph().add_run()
        run_sub.text = subline
        run_sub.font.size = Pt(16)
        run_sub.font.name = "メイリオ"

    def _add_picture_if_exists(self, slide, path: str, left, top, *, width=None, height=None) -> None:
        if path and os.path.exists(path):
            slide.shapes.add_picture(path, left, top, width=width, height=height)

    def dataset_info_to_matrix(self) -> pd.DataFrame:
        mat = []
        # 1, data count
        mat.append(
            ["データ数", self.dataset_info["Data_count"]]
        )
        # 2, Target
        mat.append(
            ["目的変数", self.dataset_info["Target_value"]["Target_name"]]
        )
        # 3, Good variable
        mat.append(["   Good sample", " "])
        mat.append(
            ["      サンプル数", self.dataset_info["Target_value"]["Good"]["Count"]]
        )
        if not pd.isna(self.dataset_info["Target_value"]["Good"]["Value_average"]):
            mat.append(
                ["      平均値", format(self.dataset_info["Target_value"]["Good"]["Value_average"], ".2e")]
            )
        # 4, Bad variable
        mat.append(["   Bad sample", " "])
        mat.append(
            ["      サンプル数", self.dataset_info["Target_value"]["Bad"]["Count"]]
        )
        if not pd.isna(self.dataset_info["Target_value"]["Bad"]["Value_average"]):
            mat.append(
                ["      平均値", format(self.dataset_info["Target_value"]["Bad"]["Value_average"], ".2e")]
            )
        # 5, Exploratory variable
        mat.append(["説明変数", " "])
        mat.append(
            ["   連続変数", self.dataset_info["Explanatory_value"]["Numeric_params_count"]]
        )
        mat.append(
            ["   カテゴリ変数", self.dataset_info["Explanatory_value"]["Categorical_params_count"]]
        )
        # dataframe
        mat = pd.DataFrame(mat, columns=["Contents", "Value"])
        return mat

    def eval_info_to_matrix(self) -> pd.DataFrame:
        mat = []
        # 1, 
        mat.append(
            ["評価データ", " "]
        )
        # 2, 評価データ数
        mat.append(
            ["  サンプル数", self.evalation_result["Validation"]["Data_count"]]
        )
        # 3, 評価データ, Accuracy
        mat.append(
            ["  Accuracy", round(self.evalation_result["Validation"]["Accuracy"],3)]
        )
        # 4, 評価データ, Recall
        mat.append(
            ["  Recall", round(self.evalation_result["Validation"]["Recall"],3)]
        )
        # 5, 評価データ, Precision
        mat.append(
            ["  Recall", round(self.evalation_result["Validation"]["Precision"],3)]
        )
        # 6, 評価データ, F1 score
        mat.append(
            ["  Recall", round(self.evalation_result["Validation"]["F1_score"],3)]
        )
        # 空欄
        mat.append(
            [" ", " "]
        )
        # 7, 
        mat.append(
            ["学習データ", " "]
        )
        # 8, 学習データ数
        mat.append(
            ["  サンプル数", self.evalation_result["Train"]["Data_count"]]
        )
        # 9, 学習データ, Accuracy
        mat.append(
            ["  Accuracy", round(self.evalation_result["Train"]["Accuracy"],3)]
        )
        # 10, 学習データ, Recall
        mat.append(
            ["  Recall", round(self.evalation_result["Train"]["Recall"],3)]
        )
        # 11, 学習データ, Precision
        mat.append(
            ["  Recall", round(self.evalation_result["Train"]["Precision"],3)]
        )
        # 12, 学習データ, F1 score
        mat.append(
            ["  Recall", round(self.evalation_result["Train"]["F1_score"],3)]
        )
        # dataframe
        mat = pd.DataFrame(mat, columns=["Contents", "Value"])
        return mat

    def importance_matrix(self) -> pd.DataFrame:
        importance_df_path = self.analyze_result["importance_df"]
        importance_df = pd.read_csv(importance_df_path)
        # sort
        importance_df.sort_values(by="importance", ascending=False)
        # ranking
        importance_df["ランキング"] = np.arange(1, len(importance_df)+1)
        # rename
        importance_df.rename(columns={"feature":"変数"}, inplace=True)
        # slice and reset index
        importance_df = importance_df[["ランキング", "変数"]]
        importance_df.reset_index(drop=True, inplace=True)
        # only top 6
        importance_df = importance_df.iloc[:6,:]

        return importance_df

    def add_box(self, slide, text, left, top, width, height, box_rgb):
        """Add a border‑less rectangle filled dark‑blue with centred white text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(left), int(top), int(width), int(height)
        )
        # remove outline
        shape.line.fill.background()
        # fill colour
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_rgb

        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)  # white text
        p.alignment = 1  # centre
        return shape

    def connect(self, parent, child, slide):
        """Draw an elbow connector from parent's right to child's left."""
        PARENT_CXN = 3   # parent : right‑center
        CHILD_CXN  = 1   # child  : left‑center
        LINE_WIDTH_PT = 0.75            # thinner connectors

        start_x = parent.left + parent.width
        start_y = parent.top + parent.height // 2
        end_x   = child.left
        end_y   = child.top + child.height // 2

        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW,
            int(start_x), int(start_y), int(end_x), int(end_y)
        )
        conn.begin_connect(parent, PARENT_CXN)
        conn.end_connect(child, CHILD_CXN)
        conn.line.width = Pt(LINE_WIDTH_PT)

    def save_pptx(self) -> pd.DataFrame:
        self.prs.save(self.savefilename)


if __name__ == "__main__":

    # savedir
    savedir = r"C:\Users\yktkk\Desktop\DS_practice\machine_learning\decision_tree_commonality\result\test3"

    # dataset
    f = r"C:\Users\yktkk\Desktop\DS_practice\machine_learning\decision_tree_commonality\dataset\house_price_dataset.csv"
    df = pd.read_csv(f, encoding="CP932")

    target_col = "SalePrice"
    df.dropna(subset=[target_col], inplace=True)
    df = df.iloc[:,-4:]

    # class
    Com = Commonality(
        savedir = savedir
    )

    # learning
    Com.learning(
        dataset = df,
        target_col = target_col,
        threshold = 250000,
        bad_direction = "under",
        numeric = True,
        val_size = 0.25
    )
    # analyze
    Com.analyze()

    # Report
    savefilename = os.path.join(savedir, "report", "learning_and_analyze_report.pptx")
    Rep = Report(savefilename=savefilename,
                 dataset_info = Com.dataset_info,
                 learning_result = Com.learning_result,
                 evalation_result = Com.evals,
                 analyze_result = Com.analyze_result,
                 branch_result = Com.branch)
    Rep.make_report()
